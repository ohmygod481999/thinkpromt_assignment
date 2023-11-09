import os
from typing import Iterator, List, Tuple
from pptx.presentation import Presentation as PS
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy
from deep_translator import GoogleTranslator


class PptxExtractor:
    def __init__(self, presentation_file_path) -> None:
        self.prs = Presentation(presentation_file_path)

    def paragraph_shapes_iter(self) -> Iterator[Shape]:
        a = iter("asdasd")
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                yield shape

    def picture_shapes_iter(self):
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    yield shape

    def save_texts_to_files(self, file: str) -> None:
        os.makedirs(os.path.dirname(file), exist_ok=True)
        with open(file, "w") as f:
            for shape in self.paragraph_shapes_iter():
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        f.write(run.text)
                        f.write("\n")


    def extract_images(self, folder_path: str) -> None:
        for picture in self.picture_shapes_iter():
            image = picture.image
            # ---get image "file" contents---
            image_bytes = image.blob
            # ---make up a name for the file, e.g. 'image.jpg'---

            image_filename = f'{folder_path}/image.{image.ext}'
            os.makedirs(os.path.dirname(image_filename), exist_ok=True)
            with open(image_filename, 'wb') as f:
                f.write(image_bytes)

    def translate(self, target_save_path):
        ps = [shape for shape in self.paragraph_shapes_iter()]
        for shape in ps:

            if shape.text_frame.text == "":
                continue
            clone_ps: List[_Paragraph] = []
            for p in shape.text_frame.paragraphs:
                clone_ps.append(copy.deepcopy(p))

            shape.text_frame.clear()
            for i, p in enumerate(clone_ps):
                new_p = shape.text_frame.add_paragraph()
                for r in p.runs:
                    new_r = new_p.add_run()
                    font = new_r.font
                    # font.color.rgb = r.font.color.rgb
                    font.bold = r.font.bold
                    font.italic = r.font.italic
                    font.size = r.font.size
                    new_r.text = r.text

                # Add translated paragraph under
                new_p = shape.text_frame.add_paragraph()
                new_p.level = 0
                for r in p.runs:
                    new_r = new_p.add_run()
                    font = new_r.font
                    font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
                    font.bold = r.font.bold
                    font.italic = r.font.italic
                    font.size = Pt(8)
                    translated_text =  GoogleTranslator(source='auto', target='vi').translate(r.text)  # output -> Weiter so, du bist großartig
                    new_r.text = translated_text


        os.makedirs(os.path.dirname(target_save_path), exist_ok=True)
        self.prs.save(target_save_path)
